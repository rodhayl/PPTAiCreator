"""Utilities to assemble a PowerPoint presentation using python‑pptx.

The builder accepts a list of `SlideContent` objects and a list of
reference strings, constructs a presentation with consistent styling, adds
slides for agenda, section dividers, and educational features, and writes
the file to disk.

Optionally, a template can be specified to apply custom designs to the
presentation content.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import List, Optional, Union

from pptx import Presentation
from pptx.util import Inches, Pt

from ..schemas import SlideContent, FormattedBullet, PresentationOutline
from ..utils import get_section_boundaries

logger = logging.getLogger(__name__)


class PresentationConfig:
    """Configuration for presentation generation."""

    def __init__(
        self,
        logo_path: Optional[str] = None,
        footer_text: Optional[str] = None,
        show_page_numbers: bool = True,
        font_name: str = "Calibri",
        title_font_size: int = 44,
        body_font_size: int = 20,
    ):
        self.logo_path = logo_path
        self.footer_text = footer_text
        self.show_page_numbers = show_page_numbers
        self.font_name = font_name
        self.title_font_size = title_font_size
        self.body_font_size = body_font_size


def build_presentation(
    slides: List[SlideContent],
    references: List[str],
    output_path: str | Path,
    template_name: Optional[str] = None,
    outline: Optional[PresentationOutline] = None,
    config: Optional[PresentationConfig] = None,
) -> None:
    """Create a PPTX file at `output_path` from slide data and references.

    Args:
        slides: List of slide content objects
        references: List of reference citations
        output_path: Path to save the presentation
        template_name: Optional template name to apply (uses TemplateManager)
        outline: Optional outline for agenda and section dividers
        config: Optional presentation configuration

    If template_name is provided, the TemplateManager will be used to apply
    the template design. Otherwise, a default plain presentation is created.
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    # Use template if specified
    if template_name:
        logger.info(f"Building presentation with template: {template_name}")
        try:
            from .template_manager import TemplateManager

            tm = TemplateManager()
            tm.apply_template(
                slides, references, template_name, output_path, outline, config
            )
            return
        except Exception as e:
            logger.error(f"Template application failed: {e}, falling back to default")
            # Continue with default generation

    # Default generation (no template)
    logger.info("Building presentation with enhanced default styling")
    prs = Presentation()

    # Set default config
    if config is None:
        config = PresentationConfig()

    # Build title slide
    if slides:
        _build_title_slide(prs, slides[0], config)

    # Build agenda slide if outline provided
    if outline and len(outline.sections) > 0:
        _build_agenda_slide(prs, outline, config)

    # Build content slides
    section_indices = get_section_boundaries(slides, outline)
    current_section = 0

    for i, slide_data in enumerate(slides):
        # Add section divider if needed
        if (
            current_section < len(section_indices)
            and i == section_indices[current_section]
        ):
            if current_section > 0:  # Skip first section divider
                _build_section_divider(prs, outline.sections[current_section], config)
            current_section += 1

        _build_content_slide(prs, slide_data, config)

    # Add references slide
    if references:
        _build_references_slide(prs, references, config)

    # Save the presentation
    prs.save(str(output_path))
    logger.info(f"Presentation saved to: {output_path}")


def _build_title_slide(
    prs: Presentation, slide: SlideContent, config: PresentationConfig
) -> None:
    """Build the title slide with optional logo."""
    layout = prs.slide_layouts[0]  # Title slide layout
    slide_obj = prs.slides.add_slide(layout)

    # Set title
    if slide_obj.shapes.title:
        slide_obj.shapes.title.text = slide.title

    # Set subtitle from first bullet or engagement hook
    subtitle = ""
    if slide.engagement_hook:
        subtitle = slide.engagement_hook
    elif slide.bullets:
        if isinstance(slide.bullets[0], str):
            subtitle = slide.bullets[0]
        else:
            subtitle = slide.bullets[0].text

    # Find subtitle placeholder
    for shape in slide_obj.placeholders:
        if hasattr(shape, "placeholder_format"):
            if shape.placeholder_format.type == 1:  # SUBTITLE
                shape.text = subtitle
                break

    # Add logo if configured
    if config.logo_path and Path(config.logo_path).exists():
        try:
            # Add logo to bottom right
            _logo = slide_obj.shapes.add_picture(
                config.logo_path,
                left=Inches(7),
                top=Inches(5),
                width=Inches(1.5),
                height=Inches(1),
            )
        except Exception as e:
            logger.warning(f"Could not add logo: {e}")

    # Add speaker notes with timing
    _add_enhanced_speaker_notes(slide_obj, slide, is_title_slide=True)


def _build_agenda_slide(
    prs: Presentation, outline: PresentationOutline, config: PresentationConfig
) -> None:
    """Build agenda/table of contents slide."""
    layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = "Agenda"

    # Build agenda content
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, section in enumerate(outline.sections, 1):
        p = body.add_paragraph()
        p.text = f"{i}. {section}"
        p.level = 0
        p.font.size = Pt(config.body_font_size)


def _build_section_divider(
    prs: Presentation, section_name: str, config: PresentationConfig
) -> None:
    """Build a section divider slide."""
    # Try to use section header layout (usually layout 2), fallback to title layout
    layout_idx = min(2, len(prs.slide_layouts) - 1) if len(prs.slide_layouts) > 2 else 0
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    # Set large section title
    if slide.shapes.title:
        slide.shapes.title.text = section_name
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(
            config.title_font_size
        )


def _build_content_slide(
    prs: Presentation, slide: SlideContent, config: PresentationConfig
) -> None:
    """Build a content slide with formatted bullets."""
    layout = prs.slide_layouts[1]  # Title and content layout
    slide_obj = prs.slides.add_slide(layout)

    # Set title
    title_shape = slide_obj.shapes.title
    title_shape.text = slide.title
    title_shape.text_frame.paragraphs[0].font.size = Pt(config.title_font_size)

    # Set content with formatted bullets
    content_placeholder = slide_obj.shapes.placeholders[1]
    _fill_formatted_bullets(content_placeholder, slide.bullets, config)

    # Add engagement hook as a callout if present
    if slide.engagement_hook:
        _add_engagement_hook(slide_obj, slide.engagement_hook)

    # Add Bloom's level indicator
    if slide.bloom_level:
        _add_bloom_indicator(slide_obj, slide.bloom_level)

    # Add formative check
    if slide.formative_check:
        _add_formative_check(slide_obj, slide.formative_check)

    # Add speaker notes
    _add_enhanced_speaker_notes(slide_obj, slide)


def _fill_formatted_bullets(
    placeholder, bullets: List[Union[str, FormattedBullet]], config: PresentationConfig
) -> None:
    """Fill a text frame with formatted bullets including nesting levels."""
    tf = placeholder.text_frame
    tf.clear()

    for bullet_obj in bullets:
        # Convert string bullets to FormattedBullet with defaults
        if isinstance(bullet_obj, str):
            bullet = FormattedBullet(text=bullet_obj)
        else:
            bullet = bullet_obj

        # Add paragraph
        if len(tf.paragraphs) == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Set text and formatting
        p.text = bullet.text
        p.level = min(bullet.level, 3)  # Cap at 3 levels
        p.font.size = Pt(
            config.body_font_size - (bullet.level * 2)
        )  # Smaller for nested

        # Apply formatting
        if bullet.bold:
            p.font.bold = True
        if bullet.italic:
            p.font.italic = True
        # Note: emphasis styling would require custom theme colors
        # For now, we skip emphasis color to avoid compatibility issues

    # Auto-fit text to prevent overflow
    _auto_fit_bullets(tf, config)


def _auto_fit_bullets(text_frame, config: PresentationConfig) -> None:
    """Automatically adjust font sizes to fit content within slide boundaries.

    This is a heuristic-based approach since python-pptx doesn't have native
    auto-fitting. We reduce font size for slides with many bullets or long text.

    Args:
        text_frame: The text frame to adjust
        config: Presentation configuration with font sizes
    """
    try:
        # Count total bullets and estimate content length
        total_bullets = len(text_frame.paragraphs)
        total_chars = sum(len(p.text) for p in text_frame.paragraphs)

        # Heuristic: if too many bullets or very long text, reduce font size
        max_bullets_for_full_size = 8
        max_chars_for_full_size = 800

        # Calculate reduction factor
        bullet_factor = max(
            0, (total_bullets - max_bullets_for_full_size) / max_bullets_for_full_size
        )
        char_factor = max(
            0, (total_chars - max_chars_for_full_size) / max_chars_for_full_size
        )
        reduction_factor = max(bullet_factor, char_factor, 0)

        # Reduce font size by up to 30%
        if reduction_factor > 0:
            max_reduction = 0.3
            reduction = min(reduction_factor * 0.5, max_reduction)

            for paragraph in text_frame.paragraphs:
                if paragraph.font.size:
                    original_size = paragraph.font.size.pt
                    # Don't reduce below 10pt
                    new_size = max(10, int(original_size * (1 - reduction)))
                    paragraph.font.size = Pt(new_size)

        # Check individual paragraphs for overflow (approximation)
        # If paragraph is very long, wrap it to multiple lines implicitly
        # by adjusting the text frame margin (python-pptx limitation)

    except Exception as e:
        # Don't fail on auto-fitting issues
        logger.debug(f"Auto-fit adjustment skipped: {e}")


def _build_references_slide(
    prs: Presentation, references: List[str], config: PresentationConfig
) -> None:
    """Build references slide."""
    layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = "References"

    # Build references list
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for ref in references:
        p = body.add_paragraph()
        p.text = ref
        p.level = 0
        p.font.size = Pt(config.body_font_size - 4)  # Smaller font for references


def _add_engagement_hook(slide_obj, hook_text: str) -> None:
    """Add engagement hook as a colored box on the slide."""
    try:
        # Add a text box for the hook (top of slide, below title)
        hook_box = slide_obj.shapes.add_textbox(
            left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(0.8)
        )
        tf = hook_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"💡 Engagement: {hook_text}"
        p.font.size = Pt(16)
        p.font.bold = True

        # Try to add background color (may not work in all python-pptx versions)
        try:
            hook_box.fill.solid()
            hook_box.fill.fore_color.rgb = RGB(255, 248, 220)  # Cornsilk color
        except Exception:
            pass  # Silently ignore if fill operations are not supported
    except Exception as e:
        logger.warning(f"Could not add engagement hook: {e}")


def _add_bloom_indicator(slide_obj, bloom_level: str) -> None:
    """Add Bloom's taxonomy level indicator."""
    try:
        # Colors for each level
        colors = {
            "remember": (255, 107, 107),  # Red
            "understand": (255, 178, 107),  # Orange
            "apply": (255, 230, 109),  # Yellow
            "analyze": (149, 225, 211),  # Teal
            "evaluate": (168, 230, 207),  # Green
            "create": (199, 206, 234),  # Purple
        }

        color = colors.get(bloom_level, (200, 200, 200))

        # Add colored box in top right
        indicator = slide_obj.shapes.add_textbox(
            left=Inches(8), top=Inches(0.5), width=Inches(1.5), height=Inches(0.5)
        )
        tf = indicator.text_frame
        p = tf.paragraphs[0]
        p.text = bloom_level.upper()
        p.font.size = Pt(12)
        p.font.bold = True

        # Try to set color (RGB might not be available in all versions)
        if RGB is not None:
            p.font.color.rgb = RGB(*color)

        try:
            indicator.fill.solid()
            if RGB is not None:
                indicator.fill.fore_color.rgb = RGB(*color)
        except Exception:
            pass  # Silently ignore if fill operations are not supported

    except Exception as e:
        logger.warning(f"Could not add Bloom indicator: {e}")


def _add_formative_check(slide_obj, check_text: str) -> None:
    """Add formative assessment check."""
    try:
        # Add at bottom of slide
        check_box = slide_obj.shapes.add_textbox(
            left=Inches(1), top=Inches(6.5), width=Inches(8), height=Inches(0.8)
        )
        tf = check_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"✓ Check for Understanding: {check_text}"
        p.font.size = Pt(14)
        p.font.italic = True

    except Exception as e:
        logger.warning(f"Could not add formative check: {e}")


def _add_enhanced_speaker_notes(
    slide_obj, slide: SlideContent, is_title_slide: bool = False
) -> None:
    """Add enhanced speaker notes with timing and transitions."""
    notes_slide = slide_obj.notes_slide
    notes_frame = notes_slide.notes_text_frame

    # Start with base speaker notes
    notes_text = slide.speaker_notes or ""

    # Add timing estimate
    if slide.time_estimate:
        if notes_text:
            notes_text += "\n\n"
        notes_text += f"⏱ Estimated time: {slide.time_estimate} minutes\n"

    # Add transition cue
    if not is_title_slide:
        if notes_text:
            notes_text += "\n"
        notes_text += "→ Transition to next slide"

    # Add pedagogical prompts if present
    if slide.active_learning_prompt:
        if notes_text:
            notes_text += "\n\n"
        notes_text += f"🎯 Active Learning: {slide.active_learning_prompt}"

    # Add Bloom's level info
    if slide.bloom_level:
        if notes_text:
            notes_text += "\n\n"
        notes_text += f"📚 Cognitive Level: {slide.bloom_level.upper()}"

    notes_frame.text = notes_text


try:
    from pptx.dml.color import RGB
except ImportError:
    # Fallback if RGB not available
    RGB = None
