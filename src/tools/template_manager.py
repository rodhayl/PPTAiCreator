"""Template management system for PowerPoint presentations.

This module provides functionality to load, validate, and apply design templates
to generated presentations. Templates are standard PowerPoint files with
placeholders that are automatically filled with generated content.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation

try:
    from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
except ImportError:
    from pptx.enum.shapes import PP_PLACEHOLDER as PP_PLACEHOLDER_TYPE
from pptx.util import Inches, Pt

from ..schemas import SlideContent, PresentationOutline
from ..tools.pptx_builder import PresentationConfig
from ..utils import get_section_boundaries

logger = logging.getLogger(__name__)


class TemplateValidationError(Exception):
    """Raised when template validation fails."""

    pass


class TemplateManager:
    """Manage PowerPoint templates for presentation generation.

    This class handles loading, validating, and applying design templates
    to generated content. Templates are standard .pptx files with slide
    layouts containing placeholders.
    """

    def __init__(self, template_dir: str = "templates"):
        """Initialize TemplateManager.

        Args:
            template_dir: Directory containing template files
        """
        self.template_dir = Path(template_dir)
        self.template_dir.mkdir(exist_ok=True, parents=True)

        # Create custom templates subdirectory
        self.custom_dir = self.template_dir / "custom"
        self.custom_dir.mkdir(exist_ok=True, parents=True)

        self._template_cache: Dict[str, Presentation] = {}

    def list_templates(self) -> List[str]:
        """List all available template names.

        Returns:
            List of template names (without .pptx extension)
        """
        templates = []

        # Built-in templates
        for file in self.template_dir.glob("*.pptx"):
            templates.append(file.stem)

        # Custom templates
        for file in self.custom_dir.glob("*.pptx"):
            templates.append(f"custom/{file.stem}")

        return sorted(templates)

    def get_template_path(self, template_name: str) -> Path:
        """Get full path to template file.

        Args:
            template_name: Name of template (with or without .pptx)

        Returns:
            Path to template file

        Raises:
            FileNotFoundError: If template doesn't exist
        """
        # Remove extension if provided
        if template_name.endswith(".pptx"):
            template_name = template_name[:-5]

        # Check for custom template
        if template_name.startswith("custom/"):
            template_path = self.template_dir / f"{template_name}.pptx"
        else:
            template_path = self.template_dir / f"{template_name}.pptx"

        if not template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")

        return template_path

    def load_template(self, template_name: str, use_cache: bool = True) -> Presentation:
        """Load a template file.

        Args:
            template_name: Name of template to load
            use_cache: Whether to use cached template if available

        Returns:
            Presentation object from template

        Raises:
            FileNotFoundError: If template doesn't exist
            TemplateValidationError: If template validation fails
        """
        # Check cache
        if use_cache and template_name in self._template_cache:
            logger.info(f"Using cached template: {template_name}")
            return self._template_cache[template_name]

        # Load from file
        template_path = self.get_template_path(template_name)
        logger.info(f"Loading template from: {template_path}")

        try:
            prs = Presentation(str(template_path))
        except Exception as e:
            raise TemplateValidationError(f"Failed to load template: {e}")

        # Validate template
        valid, errors = self.validate_template(prs)
        if not valid:
            error_msg = "Template validation failed:\n" + "\n".join(
                f"  - {e}" for e in errors
            )
            raise TemplateValidationError(error_msg)

        # Cache template
        if use_cache:
            self._template_cache[template_name] = prs

        logger.info(f"Template loaded successfully: {template_name}")
        return prs

    def validate_template(self, prs: Presentation) -> Tuple[bool, List[str]]:
        """Validate that template has required layouts and placeholders.

        Args:
            prs: Presentation object to validate

        Returns:
            Tuple of (is_valid, list_of_errors)
        """
        errors = []

        # Check minimum layouts
        if len(prs.slide_layouts) < 1:
            errors.append("Template must have at least 1 slide layout")
            return False, errors

        # Check title slide (layout 0)
        try:
            title_layout = prs.slide_layouts[0]
            if hasattr(title_layout, "shapes"):
                try:
                    # Try to access title - it's okay if it doesn't exist
                    _ = title_layout.shapes.title
                except AttributeError:
                    # No title attribute, check for title placeholder manually
                    has_title = any(
                        hasattr(shape, "placeholder_format")
                        and shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.TITLE
                        for shape in title_layout.placeholders
                    )
                    if not has_title:
                        # This is a warning, not a hard error
                        logger.warning(
                            "Layout 0 has no explicit title placeholder (will use first text shape)"
                        )
        except Exception as e:
            # Non-critical error, just log it
            logger.warning(f"Could not fully validate layout 0: {e}")

        # Check content layout (layout 1 or 2)
        content_layout_idx = min(1, len(prs.slide_layouts) - 1)
        try:
            content_layout = prs.slide_layouts[content_layout_idx]
            has_body = self._has_body_placeholder(content_layout)
            if not has_body:
                # Not a hard error, just a warning
                logger.warning(
                    f"Content layout {content_layout_idx} has no body placeholder"
                )
        except Exception as e:
            errors.append(f"Error checking content layout: {e}")

        is_valid = len(errors) == 0
        return is_valid, errors

    def _has_body_placeholder(self, layout) -> bool:
        """Check if layout has a body/content placeholder."""
        try:
            for shape in layout.placeholders:
                if hasattr(shape, "placeholder_format"):
                    ph_type = shape.placeholder_format.type
                    # Check for body or object placeholders (CONTENT doesn't exist in all versions)
                    if ph_type in (
                        PP_PLACEHOLDER_TYPE.BODY,
                        PP_PLACEHOLDER_TYPE.OBJECT,
                    ):
                        return True
        except Exception:
            pass
        return False

    def get_layout_for_content(
        self,
        prs: Presentation,
        content: SlideContent,
        is_first: bool = False,
        is_references: bool = False,
    ) -> int:
        """Determine which layout to use for given content.

        Args:
            prs: Presentation object
            content: Slide content to map
            is_first: Whether this is the first slide (title slide)
            is_references: Whether this is the references slide

        Returns:
            Layout index to use
        """
        num_layouts = len(prs.slide_layouts)

        # Title slide
        if is_first:
            return 0

        # References slide (prefer layout 5 if exists, else use content)
        if is_references:
            return (
                min(5, num_layouts - 1) if num_layouts > 5 else min(2, num_layouts - 1)
            )

        # Section header (check if slide has very few bullets)
        if len(content.bullets) <= 1 and not content.speaker_notes:
            # Use section layout if available (layout 1)
            return min(1, num_layouts - 1)

        # Regular content slide (layout 2 or 1)
        return min(2, num_layouts - 1) if num_layouts > 2 else min(1, num_layouts - 1)

    def apply_template(
        self,
        content: List[SlideContent],
        references: List[str],
        template_name: str,
        output_path: Path,
        outline: Optional[PresentationOutline] = None,
        config: Optional[PresentationConfig] = None,
    ) -> None:
        """Apply template to generated content and save.

        Args:
            content: List of slide content objects
            references: List of reference citations
            template_name: Name of template to apply
            output_path: Path to save final presentation
            outline: Optional outline for agenda and section dividers
            config: Optional presentation configuration

        Raises:
            TemplateValidationError: If template is invalid
        """
        logger.info(f"Applying template '{template_name}' to content")

        # Load template
        prs = self.load_template(template_name)

        # Set default config
        if config is None:
            config = PresentationConfig()

        # Add title slide
        if content:
            first_slide = content[0]
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            self._fill_title_slide(slide, first_slide, config)

            # Add agenda slide if outline provided
            if outline and len(outline.sections) > 0:
                agenda_layout = prs.slide_layouts[min(1, len(prs.slide_layouts) - 1)]
                agenda_slide = prs.slides.add_slide(agenda_layout)
                agenda_slide.shapes.title.text = "Agenda"
                body = agenda_slide.shapes.placeholders[1].text_frame
                body.clear()
                for i, section in enumerate(outline.sections, 1):
                    p = body.add_paragraph()
                    p.text = f"{i}. {section}"
                    p.level = 0

            # Add remaining content slides
            section_indices = get_section_boundaries(content, outline)
            current_section = 0

            for i, slide_content in enumerate(content[1:], 1):
                # Add section divider if needed
                if (
                    current_section < len(section_indices)
                    and i == section_indices[current_section]
                ):
                    if current_section > 0:
                        section_name = (
                            outline.sections[current_section]
                            if outline
                            else f"Section {current_section + 1}"
                        )
                        section_layout_idx = (
                            min(2, len(prs.slide_layouts) - 1)
                            if len(prs.slide_layouts) > 2
                            else 0
                        )
                        section_layout = prs.slide_layouts[section_layout_idx]
                        section_slide = prs.slides.add_slide(section_layout)
                        if section_slide.shapes.title:
                            section_slide.shapes.title.text = section_name
                    current_section += 1

                layout_idx = self.get_layout_for_content(
                    prs, slide_content, is_first=False
                )
                layout = prs.slide_layouts[layout_idx]
                slide = prs.slides.add_slide(layout)
                self._fill_content_slide(slide, slide_content, config)

        # Add references slide
        if references:
            layout_idx = self.get_layout_for_content(
                prs,
                SlideContent(title="References", bullets=[], speaker_notes=""),
                is_references=True,
            )
            layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(layout)
            self._fill_references_slide(slide, references, config)

        # Save presentation
        prs.save(str(output_path))
        logger.info(f"Template applied successfully: {output_path}")

    def _fill_title_slide(
        self, slide, content: SlideContent, config: PresentationConfig
    ) -> None:
        """Fill title slide with content.

        Args:
            slide: Slide object to fill
            content: Content for title slide
            config: Presentation configuration
        """
        try:
            # Fill title
            if slide.shapes.title:
                slide.shapes.title.text = content.title

            # Fill subtitle (usually first placeholder after title)
            for shape in slide.placeholders:
                if (
                    hasattr(shape, "placeholder_format")
                    and shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.SUBTITLE
                ):
                    # Use engagement hook or first bullet as subtitle
                    subtitle_text = ""
                    if content.engagement_hook:
                        subtitle_text = content.engagement_hook
                    elif content.bullets:
                        subtitle_text = (
                            content.bullets[0]
                            if isinstance(content.bullets[0], str)
                            else content.bullets[0].text
                        )
                    shape.text = subtitle_text
                    break

            # Add logo if configured
            if config.logo_path and Path(config.logo_path).exists():
                try:
                    _logo = slide.shapes.add_picture(
                        config.logo_path,
                        left=Inches(7),
                        top=Inches(5),
                        width=Inches(1.5),
                        height=Inches(1),
                    )
                except Exception as e:
                    logger.warning(f"Could not add logo: {e}")
        except Exception as e:
            logger.warning(f"Error filling title slide: {e}")

    def _fill_content_slide(
        self, slide, content: SlideContent, config: PresentationConfig
    ) -> None:
        """Fill content slide with bullets and notes.

        Args:
            slide: Slide object to fill
            content: Content for slide
            config: Presentation configuration
        """
        try:
            # Fill title
            if slide.shapes.title:
                slide.shapes.title.text = content.title

            # Fill body placeholder with bullets
            body_shape = None
            for shape in slide.placeholders:
                if hasattr(shape, "placeholder_format"):
                    ph_type = shape.placeholder_format.type
                    # Check for body or object placeholders
                    if ph_type in (
                        PP_PLACEHOLDER_TYPE.BODY,
                        PP_PLACEHOLDER_TYPE.OBJECT,
                    ):
                        body_shape = shape
                        break

            if body_shape and hasattr(body_shape, "text_frame"):
                text_frame = body_shape.text_frame
                text_frame.clear()

                # Add formatted bullets
                for i, bullet in enumerate(content.bullets):
                    if i == 0:
                        # Use first paragraph
                        p = text_frame.paragraphs[0]
                    else:
                        # Add new paragraph
                        p = text_frame.add_paragraph()

                    # Handle both string and FormattedBullet
                    if isinstance(bullet, str):
                        p.text = bullet
                        p.level = 0
                    else:
                        p.text = bullet.text
                        p.level = min(bullet.level, 3)
                        if bullet.bold:
                            p.font.bold = True
                        if bullet.italic:
                            p.font.italic = True

            # Add engagement hook as a callout if present
            if content.engagement_hook:
                try:
                    hook_box = slide.shapes.add_textbox(
                        left=Inches(1),
                        top=Inches(2),
                        width=Inches(8),
                        height=Inches(0.8),
                    )
                    tf = hook_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = f"💡 Engagement: {content.engagement_hook}"
                    p.font.size = Pt(16)
                    p.font.bold = True
                except Exception as e:
                    logger.warning(f"Could not add engagement hook: {e}")

            # Add Bloom's level indicator
            if content.bloom_level:
                try:
                    colors = {
                        "remember": (255, 107, 107),
                        "understand": (255, 178, 107),
                        "apply": (255, 230, 109),
                        "analyze": (149, 225, 211),
                        "evaluate": (168, 230, 207),
                        "create": (199, 206, 234),
                    }
                    _color = colors.get(content.bloom_level, (200, 200, 200))
                    # TODO: Apply color to indicator background

                    indicator = slide.shapes.add_textbox(
                        left=Inches(8),
                        top=Inches(0.5),
                        width=Inches(1.5),
                        height=Inches(0.5),
                    )
                    tf = indicator.text_frame
                    p = tf.paragraphs[0]
                    p.text = content.bloom_level.upper()
                    p.font.size = Pt(12)
                    p.font.bold = True
                except Exception as e:
                    logger.warning(f"Could not add Bloom indicator: {e}")

            # Add formative check
            if content.formative_check:
                try:
                    check_box = slide.shapes.add_textbox(
                        left=Inches(1),
                        top=Inches(6.5),
                        width=Inches(8),
                        height=Inches(0.8),
                    )
                    tf = check_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = f"✓ Check for Understanding: {content.formative_check}"
                    p.font.size = Pt(14)
                    p.font.italic = True
                except Exception as e:
                    logger.warning(f"Could not add formative check: {e}")

            # Add enhanced speaker notes
            self._add_enhanced_speaker_notes(slide, content)

        except Exception as e:
            logger.warning(f"Error filling content slide '{content.title}': {e}")

    def _fill_references_slide(
        self, slide, references: List[str], config: PresentationConfig
    ) -> None:
        """Fill references slide with citations.

        Args:
            slide: Slide object to fill
            references: List of reference strings
            config: Presentation configuration
        """
        try:
            # Fill title
            if slide.shapes.title:
                slide.shapes.title.text = "References"

            # Fill body with references
            body_shape = None
            for shape in slide.placeholders:
                if hasattr(shape, "placeholder_format"):
                    ph_type = shape.placeholder_format.type
                    # Check for body or object placeholders
                    if ph_type in (
                        PP_PLACEHOLDER_TYPE.BODY,
                        PP_PLACEHOLDER_TYPE.OBJECT,
                    ):
                        body_shape = shape
                        break

            if body_shape and hasattr(body_shape, "text_frame"):
                text_frame = body_shape.text_frame
                text_frame.clear()

                # Add references as bullets
                for i, ref in enumerate(references):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()

                    p.text = ref
                    p.level = 0
                    # Make font smaller for references
                    for run in p.runs:
                        run.font.size = Pt(config.body_font_size - 4 if config else 14)

        except Exception as e:
            logger.warning(f"Error filling references slide: {e}")

    def _add_enhanced_speaker_notes(self, slide_obj, content: SlideContent) -> None:
        """Add enhanced speaker notes with timing and transitions.

        Args:
            slide_obj: Slide object
            content: Slide content
        """
        notes_slide = slide_obj.notes_slide
        notes_frame = notes_slide.notes_text_frame

        # Start with base speaker notes
        notes_text = content.speaker_notes or ""

        # Add timing estimate
        if content.time_estimate:
            if notes_text:
                notes_text += "\n\n"
            notes_text += f"⏱ Estimated time: {content.time_estimate} minutes\n"

        # Add transition cue
        if notes_text:
            notes_text += "\n"
        notes_text += "→ Transition to next slide"

        # Add pedagogical prompts if present
        if content.active_learning_prompt:
            if notes_text:
                notes_text += "\n\n"
            notes_text += f"🎯 Active Learning: {content.active_learning_prompt}"

        # Add Bloom's level info
        if content.bloom_level:
            if notes_text:
                notes_text += "\n\n"
            notes_text += f"📚 Cognitive Level: {content.bloom_level.upper()}"

        notes_frame.text = notes_text

    def save_uploaded_template(
        self, file_data: bytes, filename: str
    ) -> Tuple[bool, str]:
        """Save uploaded template file and validate it.

        Args:
            file_data: Binary data of uploaded file
            filename: Original filename

        Returns:
            Tuple of (success, message)
        """
        try:
            # Save to custom directory
            output_path = self.custom_dir / filename
            with open(output_path, "wb") as f:
                f.write(file_data)

            # Validate template
            template_name = f"custom/{output_path.stem}"
            try:
                _prs = self.load_template(template_name, use_cache=False)
                return True, f"Template '{filename}' uploaded successfully"
            except TemplateValidationError as e:
                # Delete invalid template
                output_path.unlink()
                return False, f"Template validation failed: {e}"

        except Exception as e:
            return False, f"Failed to save template: {e}"

    def delete_template(self, template_name: str) -> Tuple[bool, str]:
        """Delete a custom template.

        Args:
            template_name: Name of template to delete

        Returns:
            Tuple of (success, message)
        """
        try:
            # Only allow deleting custom templates
            if not template_name.startswith("custom/"):
                return False, "Can only delete custom templates"

            template_path = self.get_template_path(template_name)
            template_path.unlink()

            # Remove from cache
            if template_name in self._template_cache:
                del self._template_cache[template_name]

            return True, f"Template '{template_name}' deleted successfully"

        except FileNotFoundError:
            return False, f"Template not found: {template_name}"
        except Exception as e:
            return False, f"Failed to delete template: {e}"
