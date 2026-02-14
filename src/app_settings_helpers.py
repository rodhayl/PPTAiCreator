"""Settings management helpers for Streamlit app.

Provides functions to:
- Fetch available models from Ollama and OpenRouter
- Save/load configuration
- Display settings UI
"""

from __future__ import annotations

from pathlib import Path
from typing import List, Tuple

import requests
import streamlit as st


def fetch_ollama_models(
    base_url: str = "http://localhost:11434",
) -> Tuple[bool, List[str], str]:
    """Fetch available models from Ollama.

    Args:
        base_url: Ollama API base URL

    Returns:
        Tuple of (success, model_list, error_message)
    """
    try:
        response = requests.get(f"{base_url}/api/tags", timeout=5)
        response.raise_for_status()
        data = response.json()

        models = []
        if "models" in data:
            for model_info in data["models"]:
                model_name = model_info.get("name", model_info.get("model", ""))
                if model_name:
                    models.append(model_name)

        return True, sorted(models), ""
    except requests.exceptions.ConnectionError:
        return (
            False,
            [],
            f"Could not connect to Ollama at {base_url}. Make sure Ollama is running.",
        )
    except requests.exceptions.Timeout:
        return False, [], "Request timed out. Ollama may be unresponsive."
    except Exception as e:
        return False, [], f"Error fetching Ollama models: {str(e)}"


def fetch_openrouter_models(api_key: str) -> Tuple[bool, List[str], str]:
    """Fetch available models from OpenRouter.

    Args:
        api_key: OpenRouter API key

    Returns:
        Tuple of (success, model_list, error_message)
    """
    if not api_key or not api_key.strip():
        return False, [], "API key is required to fetch OpenRouter models"

    try:
        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        }

        response = requests.get(
            "https://openrouter.ai/api/v1/models", headers=headers, timeout=10
        )
        response.raise_for_status()
        data = response.json()

        models = []
        if "data" in data:
            for model_info in data["data"]:
                model_id = model_info.get("id", "")
                if model_id:
                    models.append(model_id)

        return True, sorted(models), ""
    except requests.exceptions.HTTPError as e:
        if e.response.status_code == 401:
            return False, [], "Invalid API key. Please check your OpenRouter API key."
        return False, [], f"HTTP error: {e.response.status_code}"
    except requests.exceptions.Timeout:
        return False, [], "Request timed out. OpenRouter may be unresponsive."
    except Exception as e:
        return False, [], f"Error fetching OpenRouter models: {str(e)}"


def load_config_from_file(config_path: Path) -> dict:
    """Load configuration from ai_config.properties file.

    Args:
        config_path: Path to config file

    Returns:
        Dictionary of configuration key-value pairs
    """
    config = {}
    if not config_path.exists():
        return config

    with open(config_path, "r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith("#"):
                continue
            if "=" in line:
                key, value = line.split("=", 1)
                config[key.strip()] = value.strip()

    return config


def save_config_to_file(config: dict, config_path: Path) -> None:
    """Save configuration to ai_config.properties file.

    Args:
        config: Dictionary of configuration key-value pairs
        config_path: Path to config file
    """
    lines = []
    lines.append("# AI Model Configuration for PPTX Agent")
    lines.append("# This file centralizes all AI model settings for the application.")
    lines.append("")

    # Provider section
    lines.append(
        "# ============================================================================"
    )
    lines.append("# PROVIDER SELECTION")
    lines.append(
        "# ============================================================================"
    )
    lines.append(f"ai.provider={config.get('ai.provider', 'ollama')}")
    lines.append("")

    # Ollama section
    lines.append(
        "# ============================================================================"
    )
    lines.append("# OLLAMA CONFIGURATION")
    lines.append(
        "# ============================================================================"
    )
    lines.append(
        f"ollama.base_url={config.get('ollama.base_url', 'http://localhost:11434/v1')}"
    )
    lines.append(f"ollama.model={config.get('ollama.model', 'gpt-oss:20b-cloud')}")
    lines.append(
        f"ollama.fallback_model={config.get('ollama.fallback_model', 'qwen3:1.7B')}"
    )
    lines.append(f"ollama.temperature={config.get('ollama.temperature', '0.7')}")
    lines.append(f"ollama.max_tokens={config.get('ollama.max_tokens', '2048')}")
    lines.append("")

    # OpenRouter section
    lines.append(
        "# ============================================================================"
    )
    lines.append("# OPENROUTER CONFIGURATION")
    lines.append(
        "# ============================================================================"
    )
    lines.append(
        f"openrouter.base_url={config.get('openrouter.base_url', 'https://openrouter.ai/api/v1')}"
    )
    lines.append(
        f"openrouter.model={config.get('openrouter.model', 'google/gemma-2-9b-it:free')}"
    )
    lines.append(f"openrouter.api_key={config.get('openrouter.api_key', '')}")
    lines.append(
        f"openrouter.temperature={config.get('openrouter.temperature', '0.7')}"
    )
    lines.append(f"openrouter.max_tokens={config.get('openrouter.max_tokens', '2048')}")
    lines.append("")

    # Agent-specific overrides
    lines.append(
        "# ============================================================================"
    )
    lines.append("# AGENT-SPECIFIC MODEL OVERRIDES (OPTIONAL)")
    lines.append(
        "# ============================================================================"
    )
    lines.append(
        f"agent.brainstorm.temperature={config.get('agent.brainstorm.temperature', '0.8')}"
    )
    lines.append(
        f"agent.research.temperature={config.get('agent.research.temperature', '0.3')}"
    )
    lines.append(
        f"agent.content.temperature={config.get('agent.content.temperature', '0.7')}"
    )
    lines.append(f"agent.qa.temperature={config.get('agent.qa.temperature', '0.2')}")
    lines.append("")

    with open(config_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))


def _display_template_management() -> None:
    """Display template management UI section."""
    try:
        from src.tools.template_manager import TemplateManager

        tm = TemplateManager()
        templates = tm.list_templates()

        st.markdown("Manage PowerPoint templates for presentations.")

        # Display available templates
        if templates:
            st.markdown("**Available Templates:**")

            for template in templates:
                col1, col2, col3 = st.columns([3, 1, 1])

                with col1:
                    # Show template name with icon
                    icon = "📦" if not template.startswith("custom/") else "🎨"
                    st.text(f"{icon} {template}")

                with col2:
                    if st.button("👁 Preview", key=f"preview_{template}"):
                        try:
                            template_path = tm.get_template_path(template)
                            from pptx import Presentation

                            prs = Presentation(str(template_path))
                            layout_names = [
                                layout.name or f"Layout {idx}"
                                for idx, layout in enumerate(prs.slide_layouts)
                            ]

                            st.info(
                                f"Template: {template} | Layouts: {len(prs.slide_layouts)} | "
                                f"Size: {template_path.stat().st_size / 1024:.1f} KB"
                            )
                            with st.expander(f"📋 {template} details", expanded=False):
                                st.text(f"Path: {template_path}")
                                st.text("Layouts:")
                                for idx, name in enumerate(layout_names):
                                    st.text(f"  {idx}. {name}")
                        except Exception as e:
                            st.error(f"Could not preview template '{template}': {e}")

                with col3:
                    # Delete button (only for custom templates)
                    if template.startswith("custom/"):
                        if st.button("🗑 Delete", key=f"delete_{template}"):
                            success, message = tm.delete_template(template)
                            if success:
                                st.success(message)
                                st.rerun()
                            else:
                                st.error(message)
                    else:
                        st.text("")  # Empty for built-in templates
        else:
            st.info(
                "No templates available. Upload a template or run `python scripts/create_default_templates.py`"
            )

        # Upload new template
        st.markdown("---")
        st.markdown("**Upload Custom Template:**")

        uploaded_file = st.file_uploader(
            "Choose a PowerPoint file (.pptx)",
            type=["pptx"],
            help="Upload a PowerPoint template with slide layouts and placeholders",
            key="template_uploader",
        )

        if uploaded_file is not None:
            if st.button("📤 Upload Template", key="upload_template_btn"):
                with st.spinner("Uploading and validating template..."):
                    file_data = uploaded_file.read()
                    success, message = tm.save_uploaded_template(
                        file_data, uploaded_file.name
                    )

                    if success:
                        st.success(message)
                        st.rerun()
                    else:
                        st.error(message)

        # Template info
        with st.expander("ℹ️ Template Requirements"):
            st.markdown("""
            **Required Layouts:**
            - Layout 0: Title Slide (with title placeholder)
            - Layout 1+: Content Slides (with title and body placeholders)

            **How to Create Templates:**
            1. Open PowerPoint
            2. Design your master slides (View → Slide Master)
            3. Use built-in placeholders (Insert → Placeholder)
            4. Save as .pptx file
            5. Upload here

            **Tips:**
            - Use native PowerPoint placeholders, not text boxes
            - Test with small presentation first
            - Keep layouts simple for best results
            """)

    except ImportError as e:
        st.error(f"Template system not available: {e}")
    except Exception as e:
        st.error(f"Error loading templates: {e}")


def display_settings_ui() -> None:
    """Display comprehensive settings UI in Streamlit."""
    st.header("⚙️ Model Configuration Settings")
    st.markdown("Configure your AI models, providers, and fallback options.")

    # Load current config
    config_path = Path.cwd() / "ai_config.properties"
    current_config = load_config_from_file(config_path)

    # Initialize session state for models list
    if "available_models" not in st.session_state:
        st.session_state.available_models = []
    if "models_fetched" not in st.session_state:
        st.session_state.models_fetched = False

    # Provider selection
    st.subheader("1️⃣ Provider Selection")
    col1, col2 = st.columns([2, 1])

    with col1:
        provider = st.selectbox(
            "AI Provider",
            options=["ollama", "openrouter"],
            index=0 if current_config.get("ai.provider", "ollama") == "ollama" else 1,
            help="Choose between local Ollama or cloud-based OpenRouter",
        )

    with col2:
        st.markdown("**Current Provider**")
        if provider == "ollama":
            st.success("🏠 Local Ollama")
        else:
            st.info("☁️ Cloud OpenRouter")

    st.markdown("---")

    # Ollama Configuration
    if provider == "ollama":
        st.subheader("2️⃣ Ollama Configuration")

        col1, col2 = st.columns(2)

        with col1:
            ollama_url = st.text_input(
                "Ollama Base URL",
                value=current_config.get(
                    "ollama.base_url", "http://localhost:11434/v1"
                ),
                help="URL where Ollama is running",
            )

        with col2:
            if st.button("🔍 Fetch Available Models", key="fetch_ollama"):
                with st.spinner("Fetching models from Ollama..."):
                    # Remove /v1 suffix for API call
                    api_url = ollama_url.replace("/v1", "")
                    success, models, error = fetch_ollama_models(api_url)

                    if success:
                        st.session_state.available_models = models
                        st.session_state.models_fetched = True
                        st.success(f"✅ Found {len(models)} models!")
                    else:
                        st.error(f"❌ {error}")
                        st.session_state.available_models = []

        # Show available models if fetched
        if st.session_state.models_fetched and st.session_state.available_models:
            st.info(
                f"📦 Available Models: {', '.join(st.session_state.available_models[:5])}{'...' if len(st.session_state.available_models) > 5 else ''}"
            )

        col1, col2 = st.columns(2)

        with col1:
            if st.session_state.available_models:
                primary_model = st.selectbox(
                    "Primary Model",
                    options=st.session_state.available_models,
                    index=(
                        st.session_state.available_models.index(
                            current_config.get("ollama.model", "gpt-oss:20b-cloud")
                        )
                        if current_config.get("ollama.model", "gpt-oss:20b-cloud")
                        in st.session_state.available_models
                        else 0
                    ),
                    help="Main model for generation",
                )
            else:
                primary_model = st.text_input(
                    "Primary Model",
                    value=current_config.get("ollama.model", "gpt-oss:20b-cloud"),
                    help="Model name (fetch models to see available options)",
                )

        with col2:
            if st.session_state.available_models:
                fallback_model = st.selectbox(
                    "Fallback Model",
                    options=st.session_state.available_models,
                    index=(
                        st.session_state.available_models.index(
                            current_config.get("ollama.fallback_model", "qwen3:1.7B")
                        )
                        if current_config.get("ollama.fallback_model", "qwen3:1.7B")
                        in st.session_state.available_models
                        else 0
                    ),
                    help="Backup model for rate limiting",
                )
            else:
                fallback_model = st.text_input(
                    "Fallback Model",
                    value=current_config.get("ollama.fallback_model", "qwen3:1.7B"),
                    help="Smaller model used when rate limited",
                )

        col1, col2 = st.columns(2)

        with col1:
            temperature = st.slider(
                "Temperature",
                min_value=0.0,
                max_value=2.0,
                value=float(current_config.get("ollama.temperature", "0.7")),
                step=0.1,
                help="Higher = more creative, Lower = more focused",
            )

        with col2:
            max_tokens = st.number_input(
                "Max Tokens",
                min_value=256,
                max_value=8192,
                value=int(current_config.get("ollama.max_tokens", "2048")),
                step=256,
                help="Maximum response length",
            )

    # OpenRouter Configuration
    else:  # openrouter
        st.subheader("2️⃣ OpenRouter Configuration")

        api_key = st.text_input(
            "OpenRouter API Key",
            value=current_config.get("openrouter.api_key", ""),
            type="password",
            help="Get your API key from https://openrouter.ai/keys",
        )

        col1, col2 = st.columns(2)

        with col1:
            openrouter_url = st.text_input(
                "OpenRouter Base URL",
                value=current_config.get(
                    "openrouter.base_url", "https://openrouter.ai/api/v1"
                ),
                help="OpenRouter API endpoint",
            )

        with col2:
            if st.button("🔍 Fetch Available Models", key="fetch_openrouter"):
                if not api_key:
                    st.warning("⚠️ Please enter your API key first")
                else:
                    with st.spinner("Fetching models from OpenRouter..."):
                        success, models, error = fetch_openrouter_models(api_key)

                        if success:
                            st.session_state.available_models = models
                            st.session_state.models_fetched = True
                            st.success(f"✅ Found {len(models)} models!")
                        else:
                            st.error(f"❌ {error}")
                            st.session_state.available_models = []

        # Show available models if fetched
        if st.session_state.models_fetched and st.session_state.available_models:
            st.info(
                f"📦 Available Models: {', '.join(st.session_state.available_models[:5])}{'...' if len(st.session_state.available_models) > 5 else ''}"
            )

        if st.session_state.available_models:
            primary_model = st.selectbox(
                "Primary Model",
                options=st.session_state.available_models,
                index=(
                    st.session_state.available_models.index(
                        current_config.get(
                            "openrouter.model", "google/gemma-2-9b-it:free"
                        )
                    )
                    if current_config.get(
                        "openrouter.model", "google/gemma-2-9b-it:free"
                    )
                    in st.session_state.available_models
                    else 0
                ),
                help="Main model for generation",
            )
        else:
            primary_model = st.text_input(
                "Primary Model",
                value=current_config.get(
                    "openrouter.model", "google/gemma-2-9b-it:free"
                ),
                help="Model ID (fetch models to see available options)",
            )

        col1, col2 = st.columns(2)

        with col1:
            temperature = st.slider(
                "Temperature",
                min_value=0.0,
                max_value=2.0,
                value=float(current_config.get("openrouter.temperature", "0.7")),
                step=0.1,
                help="Higher = more creative, Lower = more focused",
            )

        with col2:
            max_tokens = st.number_input(
                "Max Tokens",
                min_value=256,
                max_value=8192,
                value=int(current_config.get("openrouter.max_tokens", "2048")),
                step=256,
                help="Maximum response length",
            )

    st.markdown("---")

    # Agent-specific settings
    with st.expander("🔧 Advanced: Agent-Specific Temperature Overrides"):
        st.markdown("Fine-tune temperature for each agent:")

        col1, col2 = st.columns(2)

        with col1:
            brainstorm_temp = st.slider(
                "Brainstorm Agent",
                min_value=0.0,
                max_value=2.0,
                value=float(current_config.get("agent.brainstorm.temperature", "0.8")),
                step=0.1,
                help="Higher for creative ideation",
            )

            content_temp = st.slider(
                "Content Agent",
                min_value=0.0,
                max_value=2.0,
                value=float(current_config.get("agent.content.temperature", "0.7")),
                step=0.1,
                help="Balanced for content generation",
            )

        with col2:
            research_temp = st.slider(
                "Research Agent",
                min_value=0.0,
                max_value=2.0,
                value=float(current_config.get("agent.research.temperature", "0.3")),
                step=0.1,
                help="Lower for factual accuracy",
            )

            qa_temp = st.slider(
                "QA Agent",
                min_value=0.0,
                max_value=2.0,
                value=float(current_config.get("agent.qa.temperature", "0.2")),
                step=0.1,
                help="Lowest for consistent evaluation",
            )

    st.markdown("---")

    # Save button
    col1, col2, col3 = st.columns([2, 1, 1])

    with col2:
        if st.button("💾 Save Configuration", type="primary", use_container_width=True):
            # Build new config
            new_config = {
                "ai.provider": provider,
            }

            if provider == "ollama":
                new_config.update(
                    {
                        "ollama.base_url": ollama_url,
                        "ollama.model": primary_model,
                        "ollama.fallback_model": fallback_model,
                        "ollama.temperature": str(temperature),
                        "ollama.max_tokens": str(max_tokens),
                        # Keep OpenRouter values from current config
                        "openrouter.base_url": current_config.get(
                            "openrouter.base_url", "https://openrouter.ai/api/v1"
                        ),
                        "openrouter.model": current_config.get(
                            "openrouter.model", "google/gemma-2-9b-it:free"
                        ),
                        "openrouter.api_key": current_config.get(
                            "openrouter.api_key", ""
                        ),
                        "openrouter.temperature": current_config.get(
                            "openrouter.temperature", "0.7"
                        ),
                        "openrouter.max_tokens": current_config.get(
                            "openrouter.max_tokens", "2048"
                        ),
                    }
                )
            else:  # openrouter
                new_config.update(
                    {
                        "openrouter.base_url": openrouter_url,
                        "openrouter.model": primary_model,
                        "openrouter.api_key": api_key,
                        "openrouter.temperature": str(temperature),
                        "openrouter.max_tokens": str(max_tokens),
                        # Keep Ollama values from current config
                        "ollama.base_url": current_config.get(
                            "ollama.base_url", "http://localhost:11434/v1"
                        ),
                        "ollama.model": current_config.get(
                            "ollama.model", "gpt-oss:20b-cloud"
                        ),
                        "ollama.fallback_model": current_config.get(
                            "ollama.fallback_model", "qwen3:1.7B"
                        ),
                        "ollama.temperature": current_config.get(
                            "ollama.temperature", "0.7"
                        ),
                        "ollama.max_tokens": current_config.get(
                            "ollama.max_tokens", "2048"
                        ),
                    }
                )

            # Agent overrides
            new_config.update(
                {
                    "agent.brainstorm.temperature": str(brainstorm_temp),
                    "agent.research.temperature": str(research_temp),
                    "agent.content.temperature": str(content_temp),
                    "agent.qa.temperature": str(qa_temp),
                }
            )

            # Save to file
            save_config_to_file(new_config, config_path)

            # Reload AI config
            from src.models.ai_config import reload_ai_config

            reload_ai_config()

            st.success(
                "✅ Configuration saved! Changes will take effect on next generation."
            )
            st.balloons()

    with col3:
        if st.button("🔄 Reset to Defaults", use_container_width=True):
            # Reset session state
            st.session_state.available_models = []
            st.session_state.models_fetched = False
            st.rerun()

    # Template Management Section
    st.markdown("---")
    st.subheader("5️⃣ Presentation Templates")
    _display_template_management()

    # Current configuration display
    with st.expander("📋 View Current Configuration"):
        st.json(current_config)
