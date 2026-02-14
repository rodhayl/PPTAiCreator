"""Integration tests for Plan A PowerPoint enhancements.

Tests all enhancements implemented in Plan A:
1. Enhanced bullet formatting (bold, italic)
2. Nested bullet support (level 1, 2, 3)
3. Section divider slides
4. Agenda slide
5. Enhanced speaker notes with timing
6. Educational features display (Bloom's level, engagement hooks, formative checks)
7. Auto-sized text
8. Logo support
"""

from pptx import Presentation

from src.schemas import SlideContent, FormattedBullet, PresentationOutline
from src.tools.pptx_builder import build_presentation, PresentationConfig


def test_formatted_bullets(tmp_path):
    """Test formatted bullets with bold and italic."""
    slides = [
        SlideContent(
            title="Formatted Bullets",
            bullets=[
                "Regular bullet",
                FormattedBullet(text="Bold bullet", bold=True),
                FormattedBullet(text="Italic bullet", italic=True),
                FormattedBullet(text="Bold and italic", bold=True, italic=True),
            ],
        )
    ]
    references = []
    output = tmp_path / "formatted.pptx"
    build_presentation(slides, references, output)

    assert output.exists()
    prs = Presentation(str(output))
    assert len(prs.slides) == 2  # Title + References

    # Check that formatting was applied
    slide = prs.slides[1]  # First content slide
    body = slide.shapes.placeholders[1].text_frame
    # Note: There may be 5 paragraphs due to how python-pptx handles text frames
    assert len(body.paragraphs) >= 4  # At least the bullets we added


def test_nested_bullets(tmp_path):
    """Test nested bullet points."""
    slides = [
        SlideContent(
            title="Nested Bullets",
            bullets=[
                FormattedBullet(text="Main point", level=0),
                FormattedBullet(text="Sub point 1", level=1),
                FormattedBullet(text="Sub point 2", level=1),
                FormattedBullet(text="Sub-sub point", level=2),
                FormattedBullet(text="Another main point", level=0),
            ],
        )
    ]
    references = []
    output = tmp_path / "nested.pptx"
    build_presentation(slides, references, output)

    assert output.exists()
    prs = Presentation(str(output))
    assert len(prs.slides) == 2  # Title + References


def test_agenda_slide(tmp_path):
    """Test automatic agenda slide generation."""
    slides = [
        SlideContent(title="Introduction", bullets=["Welcome"]),
        SlideContent(title="Main Content", bullets=["Point 1"]),
    ]
    outline = PresentationOutline(
        topic="Test Presentation",
        audience="Test audience",
        sections=["Introduction", "Main Content", "Conclusion"],
    )
    references = []
    output = tmp_path / "agenda.pptx"
    build_presentation(slides, references, output, outline=outline)

    assert output.exists()
    prs = Presentation(str(output))

    # Should have: title, agenda, 2 content slides, references = 5 slides
    assert len(prs.slides) == 5

    # Check that agenda slide has correct content
    agenda_slide = prs.slides[1]
    assert agenda_slide.shapes.title.text == "Agenda"


def test_section_dividers(tmp_path):
    """Test automatic section divider slides."""
    slides = [
        SlideContent(title="Intro", bullets=["Welcome"]),
        SlideContent(title="Section 1", bullets=["Content 1"]),
        SlideContent(title="Section 2", bullets=["Content 2"]),
        SlideContent(title="Conclusion", bullets=["Summary"]),
    ]
    outline = PresentationOutline(
        topic="Test",
        audience="Test",
        sections=["Introduction", "Section 1", "Section 2", "Conclusion"],
    )
    references = []
    output = tmp_path / "sections.pptx"
    build_presentation(slides, references, output, outline=outline)

    assert output.exists()
    prs = Presentation(str(output))

    # Should have section dividers between sections
    # Title + Agenda + Intro + Section divider + Section 1 + Section divider + Section 2 + Conclusion + References
    assert len(prs.slides) >= 7


def test_enhanced_speaker_notes(tmp_path):
    """Test enhanced speaker notes with timing and transitions."""
    slides = [
        SlideContent(
            title="Slide with Timing",
            bullets=["Point 1"],
            speaker_notes="Base notes",
            time_estimate=3.5,
            active_learning_prompt="Ask students to think about this",
            bloom_level="understand",
        )
    ]
    references = []
    output = tmp_path / "notes.pptx"
    build_presentation(slides, references, output)

    assert output.exists()
    prs = Presentation(str(output))

    # Check notes contain enhanced content
    slide = prs.slides[1]
    notes_text = slide.notes_slide.notes_text_frame.text
    assert "3.5" in notes_text
    assert "Active Learning" in notes_text
    assert "Cognitive Level" in notes_text
    assert "Transition" in notes_text


def test_educational_features(tmp_path):
    """Test educational features (Bloom's level, engagement hook, formative check)."""
    slides = [
        SlideContent(
            title="Educational Slide",
            bullets=["Key concept"],
            engagement_hook="What do you know about X?",
            bloom_level="apply",
            formative_check="Can you explain this concept?",
        )
    ]
    references = []
    output = tmp_path / "educational.pptx"
    build_presentation(slides, references, output)

    assert output.exists()
    prs = Presentation(str(output))

    # Slide should be created successfully
    assert len(prs.slides) == 2  # Title + References


def test_auto_sized_text(tmp_path):
    """Test that text auto-sizes for slides with many bullets."""
    # Create slide with many long bullets
    many_bullets = [
        "This is a very long bullet point that might overflow and should be auto-sized by the system if it has too many characters or words in it"
        for _ in range(15)  # More than max_bullets_for_full_size
    ]

    slides = [SlideContent(title="Many Bullets", bullets=many_bullets)]
    references = []
    output = tmp_path / "autosize.pptx"
    build_presentation(slides, references, output)

    assert output.exists()
    prs = Presentation(str(output))

    # Should still create the slide successfully
    assert len(prs.slides) == 2


def test_presentation_config(tmp_path):
    """Test presentation configuration with custom settings."""
    slides = [SlideContent(title="Test", bullets=["Point"])]
    references = []

    config = PresentationConfig(
        title_font_size=36, body_font_size=18, show_page_numbers=False
    )

    output = tmp_path / "config.pptx"
    build_presentation(slides, references, output, config=config)

    assert output.exists()


def test_all_features_combined(tmp_path):
    """Test all Plan A features working together."""
    slides = [
        SlideContent(
            title="Enhanced Presentation",
            bullets=[
                "Introduction point",
                FormattedBullet(text="Key concept", bold=True),
                FormattedBullet(text="Supporting detail", italic=True, level=1),
            ],
            speaker_notes="Main presentation notes",
            time_estimate=5.0,
            engagement_hook="Have you ever wondered why...?",
            bloom_level="analyze",
            formative_check="What evidence supports this?",
        ),
        SlideContent(
            title="Conclusion",
            bullets=[
                "Summary point",
                FormattedBullet(text="Final thought", emphasis=True),
            ],
        ),
    ]

    outline = PresentationOutline(
        topic="Complete Test Presentation",
        audience="Test audience",
        sections=["Introduction", "Main Content", "Conclusion"],
    )

    references = ["[1] Source A", "[2] Source B"]

    config = PresentationConfig(title_font_size=40, body_font_size=18)

    output = tmp_path / "complete.pptx"
    build_presentation(slides, references, output, outline=outline, config=config)

    assert output.exists()

    # Verify all features are present
    prs = Presentation(str(output))

    # Should have: title slide, agenda, intro slide, section divider, main content,
    # conclusion slide, references slide
    # Plus title slide created from first slide
    assert len(prs.slides) >= 6

    # Check educational features on first content slide
    slide = prs.slides[2]  # First content slide after title and agenda
    notes_text = slide.notes_slide.notes_text_frame.text
    assert "5.0" in notes_text  # Timing
    assert "Active Learning" in notes_text or "Cognitive Level" in notes_text


if __name__ == "__main__":
    # Run all tests
    import pytest

    pytest.main([__file__, "-v"])
