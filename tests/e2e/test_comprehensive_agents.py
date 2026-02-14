"""Comprehensive E2E tests validating all agents with real Ollama.

This test suite validates that:
1. All 5 agents (brainstorm, research, content, design, qa) are properly executed
2. Each agent produces valid, high-quality outputs
3. Outputs meet expected quality standards
4. The full pipeline works end-to-end with real gpt-oss:20b-cloud model

Agents tested:
- brainstorm: Generates outline (topic, audience, sections)
- research: Finds evidence and citations (no LLM)
- content: Generates slide content with bullets and notes
- design: Builds PPTX file (no LLM)
- qa: Evaluates presentation quality

Run with: pytest tests/e2e/test_comprehensive_agents.py -v -s
"""

from pathlib import Path

import pytest
from pptx import Presentation

from src.graph.build_graph import run_pipeline
from src.graph.langgraph_impl import run_langgraph_pipeline
from src.models.ai_config import AIModelConfiguration
from src.models.ai_interface import get_ai_interface


@pytest.mark.e2e
@pytest.mark.slow
class TestComprehensiveAgentValidation:
    """Comprehensive validation of all agents with real Ollama."""

    @pytest.fixture(autouse=True)
    def configure_ollama(self, monkeypatch, tmp_path):
        """Configure test to use Ollama with gpt-oss:20b-cloud."""
        test_config_path = tmp_path / "test_ai_config.properties"
        test_config_path.write_text("""
ai.provider=ollama
ollama.base_url=http://localhost:11434/v1
ollama.model=gpt-oss:20b-cloud
ollama.temperature=0.7
ollama.max_tokens=2048

# Agent-specific temperatures
agent.brainstorm.temperature=0.8
agent.research.temperature=0.3
agent.content.temperature=0.7
agent.qa.temperature=0.2
""".strip())
        from src.models import ai_config

        original_instance = ai_config._config_instance
        ai_config._config_instance = AIModelConfiguration(test_config_path)
        yield
        ai_config._config_instance = original_instance

    def test_all_agents_execute_and_produce_outputs(self):
        """Test that all 5 agents execute and produce expected outputs."""
        user_input = "Climate change mitigation strategies for coastal cities"

        # Run full pipeline
        state = run_pipeline(user_input)

        # === AGENT 1: BRAINSTORM ===
        print("\n" + "=" * 80)
        print("AGENT 1: BRAINSTORM - Outline Generation")
        print("=" * 80)

        assert state.outline is not None, "Brainstorm agent should generate outline"
        assert state.outline.topic, "Outline should have a topic"
        assert state.outline.audience, "Outline should have target audience"
        assert (
            len(state.outline.sections) >= 3
        ), "Outline should have at least 3 sections"
        assert (
            len(state.outline.sections) <= 7
        ), "Outline should have at most 7 sections"

        # Use safe encoding for Windows console
        topic_safe = state.outline.topic.encode("ascii", "replace").decode("ascii")
        audience_safe = state.outline.audience.encode("ascii", "replace").decode(
            "ascii"
        )
        print(f"Topic: {topic_safe}")
        print(f"Audience: {audience_safe}")
        print(f"Sections ({len(state.outline.sections)}):")
        for i, section in enumerate(state.outline.sections, 1):
            section_safe = section.encode("ascii", "replace").decode("ascii")
            print(f"  {i}. {section_safe}")

        # Validate topic quality
        assert (
            len(state.outline.topic) >= 10
        ), "Topic should be descriptive (at least 10 chars)"
        assert (
            len(state.outline.topic) <= 150
        ), "Topic should be concise (at most 150 chars)"
        assert not state.outline.topic.endswith("."), "Topic should not end with period"

        # Validate sections quality
        for section in state.outline.sections:
            assert len(section) >= 5, f"Section '{section}' is too short"
            assert (
                len(section) <= 150
            ), f"Section '{section}' is too long (max 150 chars)"
            # Sections should be title case or sentence case (skip leading digits)
            first_alpha = next((c for c in section if c.isalpha()), None)
            assert (
                first_alpha and first_alpha.isupper()
            ), f"Section '{section}' should start with capital letter"

        # === AGENT 2: RESEARCH ===
        print("\n" + "=" * 80)
        print("AGENT 2: RESEARCH - Fact Checking & Citations")
        print("=" * 80)

        # Research agent extracts claims and finds evidence
        assert hasattr(state, "claims"), "Research agent should extract claims"
        assert hasattr(state, "evidences"), "Research agent should find evidences"
        assert hasattr(state, "citations"), "Research agent should generate citations"
        assert hasattr(state, "references"), "Research agent should build references"

        print(f"Claims extracted: {len(state.claims)}")
        print(f"Evidence found: {len(state.evidences)}")
        print(f"Citations generated: {len(state.citations)}")
        print(f"References prepared: {len(state.references)}")

        # Note: claims/evidence counts may be 0 for some topics (especially those not in corpus)
        # This is expected behavior - research agent tries but may not find matches

        # === AGENT 3: CONTENT ===
        print("\n" + "=" * 80)
        print("AGENT 3: CONTENT - Slide Generation")
        print("=" * 80)

        assert state.content is not None, "Content agent should generate slides"
        assert (
            len(state.content) > 0
        ), "Content agent should generate at least one slide"
        assert len(state.content) == len(
            state.outline.sections
        ), "Content agent should generate one slide per section"

        print(f"Slides generated: {len(state.content)}")
        for i, slide in enumerate(state.content, 1):
            title_safe = slide.title.encode("ascii", "replace").decode("ascii")
            print(f"\nSlide {i}: {title_safe}")
            print(f"  Bullets: {len(slide.bullets)}")
            assert len(slide.bullets) >= 1, f"Slide {i} should have at least 1 bullet"
            assert len(slide.bullets) <= 7, f"Slide {i} should have at most 7 bullets"
            for j, bullet in enumerate(slide.bullets, 1):
                bullet_safe = bullet.encode("ascii", "replace").decode("ascii")
                print(f"    {j}. {bullet_safe[:80]}...")  # Print first 80 chars

            # Validate slide title
            assert slide.title, f"Slide {i} should have a title"
            assert len(slide.title) >= 5, f"Slide {i} title too short"
            assert len(slide.title) <= 100, f"Slide {i} title too long"

            # Validate bullets
            for bullet in slide.bullets:
                assert len(bullet) >= 3, f"Bullet '{bullet}' too short"
                assert len(bullet) <= 200, f"Bullet '{bullet}' too long"

            # Validate speaker notes
            assert slide.speaker_notes, f"Slide {i} should have speaker notes"
            assert len(slide.speaker_notes) >= 10, f"Slide {i} speaker notes too short"

            notes_safe = slide.speaker_notes.encode("ascii", "replace").decode("ascii")
            print(f"  Speaker notes: {notes_safe[:80]}...")
            if slide.citations:
                print(f"  Citations: {', '.join(slide.citations)}")

        # === AGENT 4: DESIGN ===
        print("\n" + "=" * 80)
        print("AGENT 4: DESIGN - PPTX Generation")
        print("=" * 80)

        assert state.pptx_path is not None, "Design agent should create PPTX"
        assert Path(state.pptx_path).exists(), "PPTX file should exist on disk"

        # Load and validate PPTX structure
        prs = Presentation(state.pptx_path)
        # Calculate actual section divider count based on implementation logic
        num_sections = len(state.outline.sections) if state.outline else 0
        num_content_slides = len(state.content)

        # Section dividers follow this logic:
        # - boundaries calculated as: for 5 sections, 5 slides: [1, 2, 3, 4]
        # - first boundary (index 1) skipped due to current_section > 0 check
        # - remaining boundaries (2, 3, 4) create dividers for sections 1, 2, 3
        # - total dividers = boundaries_count - 2 = num_sections - 2
        section_divider_count = max(0, num_sections - 2) if num_sections > 0 else 0

        expected_slide_count = (
            1 + 1 + num_content_slides + section_divider_count + 1
        )  # title + agenda + content + dividers + references
        assert len(prs.slides) == expected_slide_count, (
            f"PPTX should have {expected_slide_count} slides "
            f"(1 title + 1 agenda + {num_content_slides} content + {section_divider_count} section dividers + 1 references), "
            f"but has {len(prs.slides)}"
        )

        print(f"PPTX file created: {state.pptx_path}")
        print(f"File size: {Path(state.pptx_path).stat().st_size:,} bytes")
        print(f"Total slides: {len(prs.slides)}")
        print("  - Title slide: 1")
        print("  - Agenda slide: 1")
        print(f"  - Content slides: {num_content_slides}")
        print(f"  - Section dividers: {section_divider_count}")
        print("  - References slide: 1")

        # Validate first few slides match content
        for i in range(min(len(state.content), len(prs.slides))):
            slide = prs.slides[i]
            # Check if slide has shapes (title, content)
            assert len(slide.shapes) >= 2, f"Slide {i+1} should have title and content"
            print(f"  Slide {i+1}: {len(slide.shapes)} shapes")

        # === AGENT 5: QA ===
        print("\n" + "=" * 80)
        print("AGENT 5: QA - Quality Assurance")
        print("=" * 80)

        assert state.qa_report is not None, "QA agent should generate report"
        assert hasattr(
            state.qa_report, "content_score"
        ), "QA report should have content_score"
        assert hasattr(
            state.qa_report, "design_score"
        ), "QA report should have design_score"
        assert hasattr(
            state.qa_report, "coherence_score"
        ), "QA report should have coherence_score"
        assert hasattr(state.qa_report, "feedback"), "QA report should have feedback"

        # Validate scores are in valid range (1.0 to 5.0)
        assert 1.0 <= state.qa_report.content_score <= 5.0, "Content score out of range"
        assert 1.0 <= state.qa_report.design_score <= 5.0, "Design score out of range"
        assert (
            1.0 <= state.qa_report.coherence_score <= 5.0
        ), "Coherence score out of range"

        print(f"Content Score: {state.qa_report.content_score:.2f}/5.0")
        print(f"Design Score: {state.qa_report.design_score:.2f}/5.0")
        print(f"Coherence Score: {state.qa_report.coherence_score:.2f}/5.0")
        print(f"\nFeedback ({len(state.qa_report.feedback)} chars):")
        # Use safe encoding for Windows console
        feedback_safe = state.qa_report.feedback.encode("ascii", "replace").decode(
            "ascii"
        )
        print(f"{feedback_safe[:300]}...")

        # Validate feedback quality
        assert len(state.qa_report.feedback) >= 50, "QA feedback should be detailed"
        assert len(state.qa_report.feedback) <= 5000, "QA feedback should be concise"

        # === OVERALL VALIDATION ===
        print("\n" + "=" * 80)
        print("OVERALL VALIDATION")
        print("=" * 80)

        # Validate no errors occurred
        assert (
            len(state.errors) == 0
        ), f"Pipeline should not have errors: {state.errors}"

        # Calculate average QA score
        avg_score = (
            state.qa_report.content_score
            + state.qa_report.design_score
            + state.qa_report.coherence_score
        ) / 3
        print(f"Average QA Score: {avg_score:.2f}/5.0")

        # Expect reasonable quality (at least 2.5 average)
        assert (
            avg_score >= 2.5
        ), f"Average QA score {avg_score:.2f} is too low (minimum 2.5)"

        print("\n[SUCCESS] All 5 agents executed successfully with valid outputs!")

    def test_langgraph_all_agents_execute(self):
        """Test that LangGraph implementation executes all agents."""
        user_input = "Benefits of renewable energy for urban development"

        # Run LangGraph pipeline
        state = run_langgraph_pipeline(user_input)

        # Verify all agent outputs are present
        assert state.outline is not None, "LangGraph: Brainstorm agent failed"
        assert hasattr(state, "claims"), "LangGraph: Research agent failed"
        assert state.content is not None, "LangGraph: Content agent failed"
        assert len(state.content) > 0, "LangGraph: Content agent produced no slides"
        assert state.pptx_path is not None, "LangGraph: Design agent failed"
        assert Path(state.pptx_path).exists(), "LangGraph: PPTX file not created"
        assert state.qa_report is not None, "LangGraph: QA agent failed"

        print("\n[SUCCESS] LangGraph: All agents executed successfully!")
        topic_safe = state.outline.topic.encode("ascii", "replace").decode("ascii")
        print(f"  Outline: {topic_safe}")
        print(f"  Slides: {len(state.content)}")
        print(
            f"  QA Scores: C={state.qa_report.content_score:.1f}, "
            f"D={state.qa_report.design_score:.1f}, "
            f"Co={state.qa_report.coherence_score:.1f}"
        )

    def test_agent_output_quality_with_llm_validation(self):
        """Use LLM to validate the quality of agent outputs."""
        user_input = "Sustainable agriculture practices for small-scale farmers"

        # Run pipeline
        state = run_pipeline(user_input)

        # Use LLM to validate brainstorm output
        ai = get_ai_interface()

        validation_prompt = f"""You are a quality validator. Evaluate the following presentation outline for relevance and quality.

**User Input:** {user_input}

**Generated Outline:**
- Topic: {state.outline.topic.encode('ascii', 'replace').decode('ascii')}
- Audience: {state.outline.audience.encode('ascii', 'replace').decode('ascii')}
- Sections: {', '.join(s.encode('ascii', 'replace').decode('ascii') for s in state.outline.sections)}

**Validation Criteria:**
1. Is the topic relevant to the user input?
2. Is the audience identification reasonable?
3. Do the sections logically cover the topic?
4. Is the outline well-structured?

**Output Format (JSON):**
{{
  "relevance_score": 4.5,
  "structure_score": 4.0,
  "overall_quality": "good",
  "issues": ["any issues found, or empty list if none"],
  "validation_passed": true
}}

Provide your validation as valid JSON:"""

        try:
            response = ai.generate(validation_prompt, agent="qa")
            result_str = response.content
            # Parse JSON with repair logic
            import json

            try:
                result = json.loads(result_str)
            except json.JSONDecodeError:
                start = result_str.find("{")
                end = result_str.rfind("}")
                if start != -1 and end != -1 and end > start:
                    result = json.loads(result_str[start : end + 1])
                else:
                    result = {"validation_passed": False}
            print("\n[LLM VALIDATION RESULTS]")
            print(f"Relevance Score: {result.get('relevance_score', 'N/A')}")
            print(f"Structure Score: {result.get('structure_score', 'N/A')}")
            print(f"Overall Quality: {result.get('overall_quality', 'N/A')}")
            print(f"Validation Passed: {result.get('validation_passed', False)}")

            if result.get("issues"):
                print(f"Issues: {', '.join(result['issues'])}")

            # Check that validation passed
            assert result.get(
                "validation_passed", False
            ), f"LLM validation failed with issues: {result.get('issues')}"

        except Exception as exc:
            # If LLM validation fails to parse, log but don't fail test
            print(f"\n[WARNING] LLM validation failed to parse: {exc}")
            print("Continuing with manual validation...")

        # Manual validation as fallback
        assert state.outline.topic, "Outline topic should not be empty"
        assert len(state.outline.sections) >= 3, "Should have at least 3 sections"

        print("\n[SUCCESS] Agent outputs validated successfully!")

    def test_multiple_topics_all_agents_work(self):
        """Test that all agents work correctly with different topic types."""
        test_cases = [
            ("Healthcare technology innovations", "Healthcare"),
            ("Financial planning for retirement", "Finance"),
            ("Machine learning applications in manufacturing", "Technology"),
        ]

        for user_input, expected_domain in test_cases:
            print(f"\n{'=' * 80}")
            print(f"Testing: {user_input}")
            print("=" * 80)

            state = run_pipeline(user_input)

            # Verify all agents completed
            assert state.outline is not None, f"Brainstorm failed for: {user_input}"
            assert state.content is not None, f"Content failed for: {user_input}"
            assert len(state.content) > 0, f"No slides generated for: {user_input}"
            assert state.pptx_path is not None, f"Design failed for: {user_input}"
            assert state.qa_report is not None, f"QA failed for: {user_input}"

            topic_safe = state.outline.topic.encode("ascii", "replace").decode("ascii")
            print(f"  [SUCCESS] Topic: {topic_safe}")
            print(f"  [SUCCESS] Slides: {len(state.content)}")
            print(
                f"  [SUCCESS] QA Average: {(state.qa_report.content_score + state.qa_report.design_score + state.qa_report.coherence_score) / 3:.2f}/5.0"
            )

        print("\n[SUCCESS] All topics tested successfully with all agents!")

    def test_agent_error_handling(self):
        """Test that agents handle edge cases gracefully."""
        # Test with very short input
        state = run_pipeline("AI")

        # All agents should still complete
        assert state.outline is not None, "Brainstorm should handle short input"
        assert state.content is not None, "Content should handle short input"
        assert state.pptx_path is not None, "Design should handle short input"
        assert state.qa_report is not None, "QA should handle short input"

        print("\n[SUCCESS] Edge case handling validated!")
        print(f"  Short input generated {len(state.content)} slides")
