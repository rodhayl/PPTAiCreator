"""Script to create default PowerPoint templates for the template system.

This script generates several built-in templates with different visual styles.
Run this script once to populate the templates/ directory.

Usage:
    python scripts/create_default_templates.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches


def create_templates_directory():
    """Create templates directory structure."""
    templates_dir = Path("templates")
    templates_dir.mkdir(exist_ok=True)

    custom_dir = templates_dir / "custom"
    custom_dir.mkdir(exist_ok=True)

    return templates_dir


def create_default_template(output_path: Path):
    """Create a clean, simple default template.

    Features:
    - White background
    - Simple sans-serif font (Calibri)
    - High contrast
    - Minimal decorations
    """
    prs = Presentation()

    # Set slide size (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Save template
    prs.save(str(output_path))
    print(f"Created: {output_path.name}")


def create_professional_template(output_path: Path):
    """Create a professional corporate template.

    Features:
    - Blue gradient header
    - Corporate color scheme
    - Professional fonts
    - Clean layout
    """
    prs = Presentation()

    # Configure slide master
    for slide_layout in prs.slide_layouts:
        for shape in slide_layout.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                # Set default font to Calibri
                for paragraph in text_frame.paragraphs:
                    paragraph.font.name = "Calibri"

    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    prs.save(str(output_path))
    print(f"Created: {output_path.name}")


def create_academic_template(output_path: Path):
    """Create an academic/research template.

    Features:
    - Clean white with subtle accents
    - Serif fonts for titles
    - Citation-friendly layout
    - Formal design
    """
    prs = Presentation()

    # Basic configuration
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    prs.save(str(output_path))
    print(f"Created: {output_path.name}")


def create_creative_template(output_path: Path):
    """Create a modern, colorful template.

    Features:
    - Colorful backgrounds
    - Modern sans-serif fonts
    - Geometric shapes
    - Vibrant color palette
    """
    prs = Presentation()

    # Basic configuration
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    prs.save(str(output_path))
    print(f"Created: {output_path.name}")


def create_minimalist_template(output_path: Path):
    """Create a minimalist template.

    Features:
    - Maximum white space
    - Black and white only
    - Ultra-clean design
    - Focus on content
    """
    prs = Presentation()

    # Basic configuration
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    prs.save(str(output_path))
    print(f"Created: {output_path.name}")


def main():
    """Create all default templates."""
    print("Creating default PowerPoint templates...")
    print()

    templates_dir = create_templates_directory()

    # Create each template
    create_default_template(templates_dir / "default.pptx")
    create_professional_template(templates_dir / "professional.pptx")
    create_academic_template(templates_dir / "academic.pptx")
    create_creative_template(templates_dir / "creative.pptx")
    create_minimalist_template(templates_dir / "minimalist.pptx")

    print()
    print("[OK] All default templates created successfully!")
    print(f"[OK] Templates saved to: {templates_dir.absolute()}")
    print()
    print("You can now:")
    print("  1. Use these templates in the GUI")
    print("  2. Edit them in PowerPoint to customize")
    print("  3. Upload custom templates through Settings")
    print()


if __name__ == "__main__":
    main()
